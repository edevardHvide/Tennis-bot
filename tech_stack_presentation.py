"""Generate a PowerPoint presentation explaining the Availability Monitor tech stack."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
BG_DARK    = RGBColor(0x1A, 0x1A, 0x2E)
BG_MID     = RGBColor(0x16, 0x21, 0x3E)
ACCENT     = RGBColor(0x00, 0xD2, 0xFF)
ACCENT2    = RGBColor(0x7B, 0x2F, 0xF7)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
GREEN      = RGBColor(0x00, 0xE6, 0x76)
ORANGE     = RGBColor(0xFF, 0x9F, 0x43)
PINK       = RGBColor(0xFF, 0x6B, 0x9D)
YELLOW     = RGBColor(0xFF, 0xE6, 0x6D)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide_content(slide, items, left=1.5, top=2.2, width=10, font_size=20,
                              color=LIGHT_GRAY, spacing=0.55, bullet_color=ACCENT):
    for i, (title, desc) in enumerate(items):
        y = top + i * spacing
        # Bullet point
        add_text_box(slide, left - 0.4, y, 0.3, 0.4, ">", font_size=20,
                     color=bullet_color, bold=True, font_name="Consolas")
        # Title + description
        txBox = slide.shapes.add_textbox(Inches(left), Inches(y), Inches(width), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = title
        run1.font.size = Pt(font_size)
        run1.font.color.rgb = WHITE
        run1.font.bold = True
        run1.font.name = "Calibri"
        if desc:
            run2 = p.add_run()
            run2.text = f"  —  {desc}"
            run2.font.size = Pt(font_size - 2)
            run2.font.color.rgb = color
            run2.font.name = "Calibri"


def add_card(slide, left, top, width, height, title, items, accent_color=ACCENT):
    """Add a rounded-rectangle card with title and bullet items."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x20, 0x2A, 0x44)
    shape.line.fill.background()
    shape.shadow.inherit = False
    # Accent bar at top
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(0.06)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    # Title
    add_text_box(slide, left + 0.25, top + 0.15, width - 0.5, 0.5,
                 title, font_size=20, color=accent_color, bold=True)
    # Items
    for i, item in enumerate(items):
        add_text_box(slide, left + 0.25, top + 0.65 + i * 0.38, width - 0.5, 0.4,
                     item, font_size=15, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)
add_text_box(slide, 1, 1.8, 11, 1.2, "Availability Monitor", font_size=48,
             color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.0, 11, 0.8, "Tech Stack Deep Dive", font_size=32,
             color=WHITE, bold=False, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.2, 11, 0.6,
             "A serverless app that monitors tennis & padel court availability",
             font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.5, 11, 0.5, "Built with Python  |  React  |  AWS",
             font_size=16, color=ACCENT2, alignment=PP_ALIGN.CENTER, font_name="Consolas")

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — What Does the App Do?
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "What Does the App Do?", font_size=36,
             color=ACCENT, bold=True)

items = [
    ("Scrapes matchi.se", "A booking website for tennis & padel courts in Norway"),
    ("No official API", "We parse raw HTML to extract available time slots"),
    ("Monitors 11 facilities", "Across Oslo and Bergen, for both tennis and padel"),
    ("Matches user preferences", "Sport, facility, days of the week, time windows, court type"),
    ("Sends email alerts", "When new courts become available that match your criteria"),
    ("Weekly newsletter", "Summary of upcoming availability across all your preferences"),
    ("React frontend", "Users register, set preferences, and submit feature requests"),
]
add_bullet_slide_content(slide, items, top=1.8, font_size=20, spacing=0.65)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Architecture Overview
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "Architecture Overview", font_size=36,
             color=ACCENT, bold=True)
add_text_box(slide, 0.8, 1.3, 11, 0.5,
             "Serverless event-driven architecture — no servers to manage",
             font_size=18, color=LIGHT_GRAY)

# Flow diagram as text
flow_items = [
    ("1. EventBridge (Cron)", "Triggers the Scraper Lambda on a schedule"),
    ("2. Scraper Lambda", "Fetches HTML from matchi.se, parses court slots"),
    ("3. DynamoDB", "Stores current availability; scraper computes diff vs. previous"),
    ("4. Notifications Lambda", "Receives diff, matches against user preferences, sends emails"),
    ("5. API Gateway + Preferences Lambda", "REST API for the frontend to manage user settings"),
    ("6. Frontend (S3 + CloudFront)", "React app served as static files from S3"),
    ("7. Newsletter Lambda", "Weekly cron job sends availability summary emails"),
    ("8. Feedback Lambda", "Saves feature requests to DynamoDB + creates GitHub issues"),
]
add_bullet_slide_content(slide, flow_items, top=2.0, font_size=18, spacing=0.6)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Backend: Python
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "Backend — Python 3.11", font_size=36,
             color=ACCENT, bold=True)
add_text_box(slide, 0.8, 1.3, 11, 0.5,
             "Each Lambda function is a small, focused Python module",
             font_size=18, color=LIGHT_GRAY)

add_card(slide, 0.8, 2.0, 5.5, 3.0, "Key Libraries", [
    "requests  — HTTP client for fetching web pages",
    "BeautifulSoup4  — HTML parsing & scraping",
    "boto3  — AWS SDK for Python (DynamoDB, SES, Lambda)",
    "arrow  — Friendlier date/time handling than datetime",
    "jinja2  — HTML email template rendering",
], accent_color=ACCENT)

add_card(slide, 7.0, 2.0, 5.5, 3.0, "Why Python?", [
    "Great ecosystem for web scraping",
    "boto3 makes AWS integration seamless",
    "Easy to read and maintain",
    "Fast prototyping for Lambda functions",
    "Huge community & learning resources",
], accent_color=GREEN)

add_text_box(slide, 0.8, 5.3, 11, 0.8,
             'Example:  requests.get("https://matchi.se/...") -> BeautifulSoup(html) -> extract slots',
             font_size=16, color=YELLOW, font_name="Consolas")

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — Frontend: React + TypeScript
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "Frontend — React + TypeScript", font_size=36,
             color=ACCENT, bold=True)

add_card(slide, 0.8, 1.6, 3.7, 3.5, "React 19", [
    "Component-based UI library",
    "Declarative: describe what to render",
    "Virtual DOM for efficient updates",
    "Huge ecosystem of tools & libs",
    "Industry standard for web apps",
], accent_color=ACCENT)

add_card(slide, 4.9, 1.6, 3.7, 3.5, "TypeScript 5.9", [
    "JavaScript + static types",
    "Catches bugs at compile time",
    "Better IDE autocomplete",
    "Self-documenting code",
    "Types for API responses",
], accent_color=ACCENT2)

add_card(slide, 9.0, 1.6, 3.7, 3.5, "Vite + Tailwind", [
    "Vite: blazing fast dev server & bundler",
    "Hot Module Replacement (HMR)",
    "Tailwind CSS: utility-first styling",
    "No custom CSS files needed",
    "Rapid UI prototyping",
], accent_color=PINK)

add_text_box(slide, 0.8, 5.5, 11, 0.8,
             "Users register by email, pick tennis/padel, set preferences, and get alerts",
             font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — AWS Services
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "AWS Cloud Services", font_size=36,
             color=ACCENT, bold=True)
add_text_box(slide, 0.8, 1.2, 11, 0.5,
             '"Serverless" = AWS manages the servers; you only write code & pay per use',
             font_size=18, color=LIGHT_GRAY)

add_card(slide, 0.5, 1.9, 3.8, 2.5, "Compute", [
    "AWS Lambda — runs code on demand",
    "No servers to provision or manage",
    "Scales automatically (0 to thousands)",
    "Pay only when code is running",
], accent_color=ORANGE)

add_card(slide, 4.7, 1.9, 3.8, 2.5, "Database", [
    "DynamoDB — NoSQL key-value store",
    "On-demand capacity (no provisioning)",
    "Single-digit ms latency",
    "5 tables for different data types",
], accent_color=GREEN)

add_card(slide, 8.9, 1.9, 3.8, 2.5, "Networking & Delivery", [
    "API Gateway — HTTP API for frontend",
    "S3 — hosts static frontend files",
    "EventBridge — cron job scheduling",
    "SES — transactional email sending",
], accent_color=ACCENT)

add_text_box(slide, 0.8, 4.8, 11, 0.8,
             "Region: eu-north-1 (Stockholm)  —  closest AWS region to Norway",
             font_size=16, color=YELLOW, alignment=PP_ALIGN.CENTER, font_name="Consolas")

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — DynamoDB Tables
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "DynamoDB — Data Model", font_size=36,
             color=ACCENT, bold=True)
add_text_box(slide, 0.8, 1.2, 11, 0.5,
             "NoSQL: each table has a Partition Key (PK) and optional Sort Key (SK)",
             font_size=18, color=LIGHT_GRAY)

tables = [
    ("tennis-users", "userId", "—", "User registration (email)"),
    ("tennis-preferences", "userId", "preferenceId", "Sport, facility, days, time, court type"),
    ("tennis-availability", "facilityId", "date", 'PK = "facility#sport" (e.g. "ota#padel")'),
    ("tennis-notifications", "notificationId", "—", "Deduplication with 24h TTL"),
    ("tennis-feedback", "feedbackId", "—", "User feature requests"),
]

# Table header
y_start = 2.0
add_text_box(slide, 1.0, y_start, 3.0, 0.45, "Table", font_size=17, color=ACCENT, bold=True)
add_text_box(slide, 4.0, y_start, 2.5, 0.45, "Partition Key", font_size=17, color=ACCENT, bold=True)
add_text_box(slide, 6.3, y_start, 2.0, 0.45, "Sort Key", font_size=17, color=ACCENT, bold=True)
add_text_box(slide, 8.3, y_start, 4.5, 0.45, "Notes", font_size=17, color=ACCENT, bold=True)

for i, (name, pk, sk, notes) in enumerate(tables):
    y = y_start + 0.55 + i * 0.55
    c = WHITE if i % 2 == 0 else LIGHT_GRAY
    add_text_box(slide, 1.0, y, 3.0, 0.45, name, font_size=15, color=YELLOW, font_name="Consolas")
    add_text_box(slide, 4.0, y, 2.5, 0.45, pk, font_size=15, color=c, font_name="Consolas")
    add_text_box(slide, 6.3, y, 2.0, 0.45, sk, font_size=15, color=c, font_name="Consolas")
    add_text_box(slide, 8.3, y, 4.5, 0.45, notes, font_size=15, color=c)

add_text_box(slide, 0.8, 5.3, 11, 0.8,
             'Key insight: composite keys like "ota#padel" encode both facility and sport in one field',
             font_size=16, color=GREEN, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — Web Scraping
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "Web Scraping — How It Works", font_size=36,
             color=ACCENT, bold=True)

steps = [
    ("Step 1: Fetch HTML", "Send HTTP GET to matchi.se for each facility + sport combo"),
    ("Step 2: Parse HTML", "BeautifulSoup extracts the booking table from raw HTML"),
    ("Step 3: Extract Slots", "Find available time slots (date, time, court name, price)"),
    ("Step 4: Diff", "Compare new slots against the previous snapshot in DynamoDB"),
    ("Step 5: Notify", "If there are NEW slots, trigger the Notifications Lambda"),
    ("Step 6: Store", "Save the current snapshot to DynamoDB for next comparison"),
]
add_bullet_slide_content(slide, steps, top=1.8, font_size=20, spacing=0.7, bullet_color=GREEN)

add_text_box(slide, 0.8, 6.2, 11, 0.5,
             "No API = fragile! If matchi.se changes their HTML, the scraper can break.",
             font_size=16, color=ORANGE, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — Email Notifications
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "Email System", font_size=36,
             color=ACCENT, bold=True)

add_card(slide, 0.8, 1.6, 5.5, 3.2, "How Notifications Work", [
    "1. Scraper finds new available slots",
    "2. Notifications Lambda loads all user preferences",
    "3. Matcher checks: facility + sport + day + time + court type",
    "4. Dedup prevents sending same alert twice (24h window)",
    "5. Jinja2 renders a beautiful HTML email",
    "6. Email sent via AWS SES or Gmail SMTP",
], accent_color=ORANGE)

add_card(slide, 7.0, 1.6, 5.5, 3.2, "Email Delivery Options", [
    "AWS SES — scalable, cheap, production email",
    "Gmail SMTP — fallback / development option",
    "SMTP_HOST env var controls which is used",
    "HTML templates in email_templates/ folder",
    "Jinja2 templating for dynamic content",
    "Branded design with facility & sport info",
], accent_color=PINK)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — Project Structure
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "Project Structure", font_size=36,
             color=ACCENT, bold=True)

left_items = [
    "facilities.py         — shared facility config",
    "lambdas/scraper/      — web scraping logic",
    "lambdas/preferences/  — user preferences API",
    "lambdas/notifications/ — alert matching & sending",
    "lambdas/newsletter/   — weekly summary emails",
    "lambdas/feedback/     — feature request handler",
]
right_items = [
    "frontend/src/         — React + TypeScript app",
    "tests/                — pytest test suite",
    "tests/fixtures/       — HTML fixtures for testing",
    "email_templates/      — Jinja2 HTML templates",
    "scripts/              — DB migration scripts",
    "infra/                — DynamoDB & API config",
]

for i, item in enumerate(left_items):
    add_text_box(slide, 1.0, 1.6 + i * 0.5, 5.5, 0.45, item,
                 font_size=15, color=LIGHT_GRAY, font_name="Consolas")
for i, item in enumerate(right_items):
    add_text_box(slide, 7.0, 1.6 + i * 0.5, 5.5, 0.45, item,
                 font_size=15, color=LIGHT_GRAY, font_name="Consolas")

add_text_box(slide, 0.8, 5.2, 11, 0.8,
             "Key pattern: facilities.py is shared — copied into each Lambda at build time",
             font_size=16, color=GREEN, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — DevOps & Deployment
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "DevOps & Deployment", font_size=36,
             color=ACCENT, bold=True)

add_card(slide, 0.5, 1.5, 3.8, 3.5, "Build Process", [
    "Makefile orchestrates deployments",
    "Each Lambda packaged as a .zip",
    "Dependencies installed to package/",
    "Handler + shared files added on top",
    "Uploaded to AWS via CLI",
], accent_color=ACCENT)

add_card(slide, 4.7, 1.5, 3.8, 3.5, "Testing", [
    "pytest for all backend tests",
    "HTML fixtures simulate matchi.se",
    "E2E pipeline tests (scrape->notify)",
    "Scoped testing: only test what changed",
    "Frontend: TypeScript type checking",
], accent_color=GREEN)

add_card(slide, 8.9, 1.5, 3.8, 3.5, "Infrastructure", [
    "All infra is on AWS (no Terraform yet)",
    "DynamoDB tables: on-demand capacity",
    "API Gateway: HTTP API (not REST)",
    "EventBridge: cron-based scheduling",
    "S3: static site hosting for frontend",
], accent_color=ORANGE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — Key Takeaways
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, 0.8, 0.5, 11, 0.8, "Key Takeaways", font_size=36,
             color=ACCENT, bold=True)

takeaways = [
    ("Serverless is powerful", "No servers to manage — focus on business logic, not infrastructure"),
    ("Python + BeautifulSoup", "Great combo for web scraping when there's no API available"),
    ("React + TypeScript", "Industry-standard frontend stack with type safety"),
    ("Event-driven design", "Lambdas trigger each other — loosely coupled, easy to extend"),
    ("DynamoDB composite keys", 'Clever key design (e.g. "ota#padel") avoids extra tables'),
    ("NoSQL trade-offs", "Flexible schema, but you must design your access patterns upfront"),
    ("Scraping is fragile", "HTML changes can break everything — always have good tests!"),
]
add_bullet_slide_content(slide, takeaways, top=1.7, font_size=20, spacing=0.68, bullet_color=GREEN)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — Questions
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_text_box(slide, 1, 2.5, 11, 1.2, "Questions?", font_size=52,
             color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.0, 11, 0.6,
             "github.com/edevardHvide/Tennis-bot",
             font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER, font_name="Consolas")

# ── Save ──
output_path = "/home/user/Tennis-bot/Tech_Stack_Presentation.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
